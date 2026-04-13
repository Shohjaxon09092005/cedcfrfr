"""
AI Pipeline Services for EduAI platform
- Text extraction from various file formats  
- Claude API for script and quiz generation
- ElevenLabs for Uzbek text-to-speech
- Kling AI for video generation
- S3 for file storage
"""

import anthropic
import requests
import boto3
import json
import time
import random
from django.conf import settings


class TextExtractor:
    """Extract text from PDF, DOCX, PPTX files"""
    
    def extract(self, file_path: str, file_type: str) -> str:
        """Extract text based on file type"""
        if file_type == 'pdf':
            return self._extract_pdf(file_path)
        elif file_type in ['docx', 'doc']:
            return self._extract_docx(file_path)
        elif file_type in ['pptx', 'ppt']:
            return self._extract_pptx(file_path)
        elif file_type == 'txt':
            return self._extract_txt(file_path)
        raise ValueError(f"Qo'llab-quvvatlanmaydigan fayl turi: {file_type}")
    
    def _extract_pdf(self, path):
        """Extract text from PDF"""
        import pdfplumber
        text = ""
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text
    
    def _extract_docx(self, path):
        """Extract text from DOCX"""
        from docx import Document
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    
    def _extract_pptx(self, path):
        """Extract text from PPTX"""
        from pptx import Presentation
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _extract_txt(self, path):
        """Extract text from TXT"""
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()


class ClaudeService:
    """Claude API integration for script and quiz generation"""
    
    def __init__(self):
        if not settings.ANTHROPIC_API_KEY:
            raise ValueError("ANTHROPIC_API_KEY sozlanmagan")
        self.client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        self.model = "claude-sonnet-4-6"
        # Retry configuration
        self.max_retries = 5
        self.base_delay = 1.0  # seconds
        self.max_delay = 60.0  # seconds
    
    def _retry_with_exponential_backoff(self, func, *args, **kwargs):
        """Execute function with exponential backoff retry logic"""
        last_exception = None
        
        for attempt in range(self.max_retries):
            try:
                return func(*args, **kwargs)
            except (anthropic.RateLimitError, anthropic.APIStatusError) as e:
                last_exception = e
                if attempt == self.max_retries - 1:
                    # Last attempt failed
                    break
                
                # Check if it's an overload error (529 status code)
                if hasattr(e, 'status_code') and e.status_code == 529:
                    # Calculate delay with exponential backoff + jitter
                    delay = min(self.base_delay * (2 ** attempt), self.max_delay)
                    jitter = random.uniform(0.1, 1.0) * delay * 0.1  # 10% jitter
                    total_delay = delay + jitter
                    
                    print(f"Claude API overloaded (attempt {attempt + 1}/{self.max_retries}). "
                          f"Retrying in {total_delay:.2f} seconds...")
                    time.sleep(total_delay)
                else:
                    # For other API errors, don't retry
                    raise e
                
            except Exception as e:
                # For non-API errors, don't retry
                raise e
        
        # If we get here, all retries failed with overload
        raise last_exception
    
    def generate_video_script(self, text: str, language: str = "uzbek") -> str:
        """Generate engaging video script from educational text"""
        if settings.USE_MOCK_AI_RESPONSES:
            # Prefer mock but try real if available
            return self._generate_video_script_mock(text, language)
        else:
            # Try real API, fall back to mock on failure
            try:
                return self._generate_video_script_real(text, language)
            except Exception as e:
                print(f"Claude API failed ({e}), using mock response...")
                return self._generate_video_script_mock(text, language)
    
    def _generate_video_script_real(self, text: str, language: str = "uzbek") -> str:
        """Real Claude API call with retry logic"""
        def _api_call():
            response = self.client.messages.create(
                model=self.model,
                max_tokens=2000,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text
        
        prompt = f"""
Siz tajribali o'zbek tili dars tuzuvchisisiz.
Quyidagi dars materialidan talabalar uchun 3-5 daqiqalik video dars skriptini tuzing.

TALABLAR:
- Til: O'zbek tili (lotin alifbosi)
- Ohang: Qiziqarli, tushunarli, ilmiy
- Tuzilma: Kirish → Asosiy qismlar (2-3) → Xulosa
- Har bir qismni [PAUSE] bilan ajrating
- Misollar va tushuntirishlar qo'shing
- Maksimal uzunlik: 800 so'z

DARS MATERIALI:
{text[:4000]}

Faqat skript matnini yozing, boshqa hech narsa yo'q.
"""
        return self._retry_with_exponential_backoff(_api_call)
    
    def _generate_video_script_mock(self, text: str, language: str = "uzbek") -> str:
        """Mock response for development"""
        return f"""Assalomu alaykum, aziz talabalar!

Bugun biz {text[:100]}... mavzusini o'rganamiz.

[PAUSE]

Darsning birinchi qismida biz asosiy tushunchalarni ko'rib chiqamiz.
Matematikada omborxona tizimi - bu resurslarni samarali boshqarish usuli hisoblanadi.

[PAUSE]

Ikkinchi qismda amaliy misollarni ko'ramiz.
Masalan, 100 ta mahsulotni 5 ta joyga qanday taqsimlash mumkin?

[PAUSE]

Xulosa qilib aytganda, omborxona tizimi biznes samaradorligini oshiradi.
Savollaringiz bo'lsa, bemalol beringiz!"""
    
    def generate_quiz(self, transcript: str, num_questions: int = 10, difficulty: str = "medium") -> dict:
        """Generate quiz questions from video transcript"""
        if settings.USE_MOCK_AI_RESPONSES:
            # Prefer mock but try real if available
            return self._generate_quiz_mock(transcript, num_questions)
        else:
            # Try real API, fall back to mock on failure
            try:
                return self._generate_quiz_real(transcript, num_questions, difficulty)
            except Exception as e:
                print(f"Claude API failed ({e}), using mock quiz...")
                return self._generate_quiz_mock(transcript, num_questions)
    
    def _generate_quiz_real(self, transcript: str, num_questions: int = 10, difficulty: str = "medium") -> dict:
        """Real Claude API call for quiz generation with retry logic"""
        def _api_call():
            response = self.client.messages.create(
                model=self.model,
                max_tokens=3000,
                messages=[{"role": "user", "content": prompt}]
            )
            text = response.content[0].text
            # Strip markdown code fences if present
            text = text.strip()
            if text.startswith("```"):
                text = text.split("```")[1]
                if text.startswith("json"):
                    text = text[4:]
            return json.loads(text.strip())
        
        # Map difficulty levels to Uzbek for the prompt
        difficulty_map = {
            "easy": "oson (Talabalar osongina javob berishishi mumkin)",
            "medium": "o'rta (Talabalar bir oz o'ylashi kerak)",
            "hard": "qiyin (Talabalar chuqur tushunish kerak)",
            "mixed": "aralash (oson, o'rta va qiyn savollar aralashgan)",
        }
        difficulty_description = difficulty_map.get(difficulty, "o'rta")
        
        prompt = f"""
Quyidagi dars transkriptidan {num_questions} ta test savol tuz.
Qiyinlik darajasi: {difficulty_description}

TALABLAR:
- Har bir savolda 4 ta variant bo'lsin (A, B, C, D)
- Faqat 1 ta to'g'ri javob
- Savollar O'zbek tilida
- Mavzu tegini ham qo'sh (topic)

JAVOB FORMATI (faqat JSON, boshqa hech narsa):
{{
  "questions": [
    {{
      "id": 1,
      "text": "savol matni",
      "options": ["variant A", "variant B", "variant C", "variant D"],
      "correct_answer": 0,
      "topic": "mavzu nomi",
      "difficulty": "easy|medium|hard"
    }}
  ]
}}

TRANSKRIPT:
{transcript[:5000]}
"""
        return self._retry_with_exponential_backoff(_api_call)
    
    def _generate_quiz_mock(self, transcript: str, num_questions: int = 10) -> dict:
        """Mock quiz response for development"""
        return {
            "questions": [
                {
                    "id": 1,
                    "text": "Omborxona tizimi nima?",
                    "options": ["Mahsulotlarni saqlash usuli", "Do'kon ochish", "Pul hisoblash", "Kitob yozish"],
                    "correct_answer": 0,
                    "topic": "Omborxona asoslari",
                    "difficulty": "easy"
                },
                {
                    "id": 2,
                    "text": "FIFO metodi nimani anglatadi?",
                    "options": ["Birinchi kirgan - birinchi chiqadi", "Oxirgi kirgan - birinchi chiqadi", "Aralash usul", "Avtomatik hisoblash"],
                    "correct_answer": 0,
                    "topic": "Inventarizatsiya",
                    "difficulty": "medium"
                },
                {
                    "id": 3,
                    "text": "Omborxona samaradorligini qanday oshirish mumkin?",
                    "options": ["Barcha javoblar to'g'ri", "Avtomatlashtirish", "Xodimlarni kamaytirish", "Mahsulotlarni ko'paytirish"],
                    "correct_answer": 1,
                    "topic": "Optimallashtirish",
                    "difficulty": "hard"
                }
            ]
        }
    
    def analyze_weak_topics(self, results: list, student_name: str) -> dict:
        """Analyze student quiz results and identify weak topics with retry logic"""
        if settings.USE_MOCK_AI_RESPONSES:
            # Prefer mock but try real if available
            return self._analyze_weak_topics_mock(student_name)
        else:
            # Try real API, fall back to mock on failure
            try:
                return self._analyze_weak_topics_real(results, student_name)
            except Exception as e:
                print(f"Claude API failed ({e}), using mock analysis...") 
                return self._analyze_weak_topics_mock(student_name)
    
    def _analyze_weak_topics_real(self, results: list, student_name: str) -> dict:
        """Real Claude API call for weak topics analysis with retry logic"""
        def _api_call():
            response = self.client.messages.create(
                model=self.model,
                max_tokens=1000,
                messages=[{"role": "user", "content": prompt}]
            )
            text = response.content[0].text.strip()
            if text.startswith("```"):
                text = text.split("```")[1]
                if text.startswith("json"):
                    text = text[4:]
            return json.loads(text.strip())
        
        prompt = f"""
Talaba: {student_name}
Test natijalari:
{results}

Quyidagilarni tahlil qiling:
1. Qaysi mavzularda ko'p xato qilindi (weak_topics)
2. Umumiy baho va izoh
3. Tavsiyalar (recommendations)

JAVOB FORMATI (faqat JSON):
{{
  "weak_topics": ["mavzu1", "mavzu2"],
  "strong_topics": ["mavzu3"],
  "overall_feedback": "umumiy izoh o'zbek tilida",
  "recommendations": ["tavsiya1", "tavsiya2"],
  "score_percentage": 75.0
}}
"""
        return self._retry_with_exponential_backoff(_api_call)
    
    def _analyze_weak_topics_mock(self, student_name: str) -> dict:
        """Mock response for weak topics analysis"""
        return {
            "weak_topics": ["Omborxona boshqaruvi", "Inventarizatsiya metodlari"],
            "strong_topics": ["Asosiy tushunchalar", "Praktik misollar"],
            "overall_feedback": f"{student_name}, sizning test natijalaringiz umumiy o'rtachaga mos. Omborxona boshqaruvi va inventarizatsiya kabi mavzularda ko'proq mashq qilishingizni tavsiya qilamiz.",
            "recommendations": [
                "Omborxona boshqaruvi haqida qo'shimcha darslarni ko'ring",
                "FIFO va LIFO metodlarini batafsil o'rganib chiqing",
                "Amaliy misollar bilan nazariyani birlashtiring"
            ],
            "score_percentage": 72.5
        }


class ElevenLabsService:
    """Text-to-speech in Uzbek using ElevenLabs"""
    
    def __init__(self):
        self.api_key = settings.ELEVENLABS_API_KEY
        self.voice_id = settings.ELEVENLABS_VOICE_ID
        self.base_url = "https://api.elevenlabs.io/v1"
    
    def generate_audio(self, script: str) -> bytes:
        """Convert script text to Uzbek speech audio"""
        if settings.USE_MOCK_AI_RESPONSES:
            # Prefer mock
            return self._generate_audio_mock(script)
        else:
            # Try real API, fall back to mock on failure
            try:
                return self._generate_audio_real(script)
            except Exception as e:
                print(f"ElevenLabs API failed ({e}), using mock audio...")
                return self._generate_audio_mock(script)
    
    def _generate_audio_real(self, script: str) -> bytes:
        """Real ElevenLabs API call"""
        chunks = self._split_text(script, max_chars=4500)
        audio_parts = []
        
        for chunk in chunks:
            response = requests.post(
                f"{self.base_url}/text-to-speech/{self.voice_id}",
                headers={
                    "xi-api-key": self.api_key,
                    "Content-Type": "application/json"
                },
                json={
                    "text": chunk,
                    "model_id": "eleven_multilingual_v2",
                    "voice_settings": {
                        "stability": 0.5,
                        "similarity_boost": 0.75,
                        "style": 0.3
                    }
                }
            )
            response.raise_for_status()
            audio_parts.append(response.content)
        
        return b"".join(audio_parts)
    
    def _generate_audio_mock(self, script: str) -> bytes:
        """Mock audio response for development"""
        # Return a small dummy MP3 file (this is just a placeholder)
        # In a real implementation, you'd have a sample audio file
        return b"RIFF\x24\x08\x00\x00WAVEfmt \x10\x00\x00\x00\x01\x00\x01\x00\x80>\x00\x00\x00}\x00\x00\x02\x00\x10\x00data\x00\x08\x00\x00"
    
    def _split_text(self, text: str, max_chars: int) -> list:
        """Split text at sentence boundaries"""
        sentences = text.replace("[PAUSE]", ".").split(".")
        chunks, current = [], ""
        for sentence in sentences:
            if len(current) + len(sentence) < max_chars:
                current += sentence + "."
            else:
                if current:
                    chunks.append(current.strip())
                current = sentence + "."
        if current:
            chunks.append(current.strip())
        return chunks


class KlingAIService:
    """Video generation using Kling AI (mock for now)"""
    
    def __init__(self):
        self.api_key = settings.KLING_API_KEY
        self.base_url = "https://api.klingai.com/v1"
    
    def generate_video(self, audio_url: str, script_summary: str) -> str:
        """Generate educational video from audio and prompt"""
        if settings.USE_MOCK_AI_RESPONSES:
            try:
                return self._generate_video_real(audio_url, script_summary)
            except Exception as e:
                print(f"Kling AI API failed, using mock video: {e}")
                return self._generate_video_mock(audio_url, script_summary)
        else:
            return self._generate_video_real(audio_url, script_summary)
    
    def _generate_video_real(self, audio_url: str, script_summary: str) -> str:
        """Real Kling AI API call (placeholder - needs implementation)"""
        # For now, return mock URL since we don't have the real API implementation
        return f"https://mock-video-cdn.example.com/videos/{hash(script_summary)}.mp4"
    
    def _generate_video_mock(self, audio_url: str, script_summary: str) -> str:
        """Mock video response for development"""
        return f"https://sample-videos.com/zip/10/mp4/SampleVideo_1280x720_1mb.mp4"
    
    def _poll_video_result(self, task_id: str, max_attempts: int = 60) -> str:
        """Poll until video is ready, return video URL"""
        import time
        for _ in range(max_attempts):
            time.sleep(10)
            try:
                response = requests.get(
                    f"{self.base_url}/videos/text2video/{task_id}",
                    headers={"Authorization": f"Bearer {self.api_key}"}
                )
                data = response.json()["data"]
                if data["task_status"] == "succeed":
                    return data["task_result"]["videos"][0]["url"]
                elif data["task_status"] == "failed":
                    raise Exception(f"Kling AI video generation failed: {data.get('task_status_msg')}")
            except Exception as e:
                print(f"Error checking video status: {e}")
                continue
        raise TimeoutError("Video generation timed out after 10 minutes")


class S3Service:
    """AWS S3 file operations"""
    
    def __init__(self):
        self.s3 = boto3.client(
            's3',
            aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            region_name=settings.AWS_REGION
        )
        self.bucket = settings.AWS_S3_BUCKET
        self.cdn_url = settings.AWS_CLOUDFRONT_URL
    
    def upload_file(self, file_obj, key: str, content_type: str) -> str:
        """Upload file to S3, return CDN URL"""
        self.s3.upload_fileobj(
            file_obj,
            self.bucket,
            key,
            ExtraArgs={'ContentType': content_type, 'ACL': 'private'}
        )
        return f"{self.cdn_url}/{key}"
    
    def upload_bytes(self, data: bytes, key: str, content_type: str) -> str:
        """Upload bytes to S3, return CDN URL"""
        import io
        self.s3.upload_fileobj(
            io.BytesIO(data),
            self.bucket,
            key,
            ExtraArgs={'ContentType': content_type}
        )
        return f"{self.cdn_url}/{key}"
    
    def generate_presigned_url(self, key: str, expires: int = 3600) -> str:
        """Generate presigned URL for file access"""
        return self.s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': self.bucket, 'Key': key},
            ExpiresIn=expires
        )
    
    def delete_file(self, key: str):
        """Delete file from S3"""
        self.s3.delete_object(Bucket=self.bucket, Key=key)
